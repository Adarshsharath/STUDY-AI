import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import io

def extract_text_from_file(file, file_type):
    """
    Extract text from uploaded files based on file type.
    
    Args:
        file: FileStorage object from Flask request
        file_type: String indicating file extension (pdf, docx, pptx, txt)
    
    Returns:
        String containing extracted text
    """
    
    if file_type == 'pdf':
        return extract_from_pdf(file)
    elif file_type == 'docx':
        return extract_from_docx(file)
    elif file_type == 'pptx':
        return extract_from_pptx(file)
    elif file_type == 'txt':
        return extract_from_txt(file)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

def extract_from_pdf(file):
    """Extract text from PDF using pdfplumber."""
    text = []
    
    try:
        # Seek to beginning to ensure we read from start
        file.seek(0)
        file_bytes = file.read()
        
        if not file_bytes:
            print("Error: PDF file is empty or could not be read")
            return ""
            
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        
        extracted = '\n\n'.join(text)
        if not extracted.strip():
            print("Warning: No text extracted from PDF. It might be an image-only PDF.")
            
        return extracted
    except Exception as e:
        print(f"PDF extraction error: {str(e)}")
        # If it fails, return empty string so the caller can handle it
        return ""

def extract_from_docx(file):
    """Extract text from DOCX using python-docx."""
    doc = DocxDocument(io.BytesIO(file.read()))
    text = []
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text.append(paragraph.text)
    
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text.append(cell.text)
    
    return '\n\n'.join(text)

def extract_from_pptx(file):
    """Extract text from PPTX using python-pptx."""
    prs = Presentation(io.BytesIO(file.read()))
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
    
    return '\n\n'.join(text)

def extract_from_txt(file):
    """Extract text from TXT file."""
    try:
        # Try UTF-8 first
        content = file.read().decode('utf-8')
    except UnicodeDecodeError:
        # Fallback to latin-1 if UTF-8 fails
        file.seek(0)
        content = file.read().decode('latin-1')
    
    return content
